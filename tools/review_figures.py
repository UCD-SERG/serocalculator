from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt


SOURCE = Path(
    "/Users/ezramorrison/Library/CloudStorage/OneDrive-UCDavisHealth/Work/"
    "Projects/SeroEpi/Serocalculator/Publications/revision 1/"
    "Figures_v3_resubmit.pptx"
)
OUTPUT = Path(
    "/Users/ezramorrison/Documents/GitHub/serocalculator/reviewed-documents/"
    "Figures_v3_resubmit_REVIEWED.pptx"
)


def replace_run_text(shape, old: str, new: str) -> None:
    matches = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old in run.text:
                matches.append(run)
    if len(matches) != 1:
        raise ValueError(f"Expected one run containing {old!r}; found {len(matches)}")
    matches[0].text = matches[0].text.replace(old, new)


def replace_shape_text(shape, text: str, font_size: float | None = None,
                       font_name: str | None = None) -> None:
    shape.text = text
    if font_size is not None:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                if font_name is not None:
                    run.font.name = font_name


def add_equation_overlay(slide, left: int, top: int, width: int, height: int,
                         text: str, font_size: float) -> None:
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    background.line.fill.background()

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Cambria Math"
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_text_overlay(slide, left: int, top: int, width: int, height: int,
                     text: str, font_size: float, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Times New Roman"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)


def revise_slide_one(prs: Presentation) -> None:
    slide = prs.slides[0]
    shapes = slide.shapes

    replace_shape_text(
        shapes[27],
        "Defines additive cross-reactivity or non-specific binding, relative assay "
        "error, and assay detection limits.",
        11,
    )
    replace_shape_text(
        shapes[30],
        "Seroincidence is estimated from the three inputs using functions for overall "
        "or stratified analyses.",
        11,
    )
    replace_shape_text(
        shapes[70],
        "Assuming conditional independence, multiple biomarkers can contribute to one "
        "seroincidence estimate.",
        10.5,
    )
    for paragraph in shapes[70].text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 6_180_000, 3_615_000, 2_650_000, 1_520_000
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    panel.line.fill.background()
    add_text_overlay(
        slide, 6_225_000, 3_650_000, 2_550_000, 235_000,
        "Individual likelihood:", 11, bold=True
    )
    add_text_overlay(
        slide, 6_225_000, 3_875_000, 2_550_000, 300_000,
        "Integrates over possible infection times and the never-infected state.", 9.5
    )
    add_equation_overlay(
        slide,
        6_400_000,
        4_155_000,
        2_200_000,
        500_000,
        "Lᵢ(λ) = ∫₀ᵃ p(Yᵢ | T=t) λe^(−λt)dt\n"
        "+ p(Yᵢ | T=NA)e^(−λa)",
        9.5,
    )
    add_text_overlay(
        slide, 6_225_000, 4_665_000, 2_550_000, 225_000,
        "Population log-likelihood:", 11, bold=True
    )
    add_equation_overlay(
        slide,
        6_500_000,
        4_875_000,
        2_180_000,
        240_000,
        "log L(λ) = Σᵢ₌₁ⁿ log Lᵢ(λ)",
        12,
    )


def revise_slide_two(prs: Presentation) -> None:
    slide = prs.slides[1]
    shapes = slide.shapes

    replace_run_text(shapes[4], "download//n6cp3", "download/n6cp3")
    replace_run_text(
        shapes[4],
        'noise <- url("https://osf.io/download//hqy4v/") |> ',
        'noise <- "https://osf.io/download/hqy4v/" |> ',
    )
    replace_run_text(shapes[11], "# Estimate seroconversion rate", "# Estimate seroincidence rates")
    assignment_run = shapes[11].text_frame.paragraphs[1].runs[1]
    if assignment_run.text != " = ":
        raise ValueError("Unexpected assignment run in estimation code")
    assignment_run.text = " <- "
    replace_run_text(shapes[11], "dataframe", "data frame")
    replace_run_text(shapes[12], "# Plot seroconversion rates", "# Plot seroincidence rates")

    replace_shape_text(
        shapes[10],
        "b) Visualize antibody distributions & seroresponse parameters",
        14,
        "Arial",
    )
    shapes[10].width = 5_800_000
    shapes[10].height = 420_000

    replace_shape_text(
        shapes[9],
        "c) Summarize stratified seroincidence estimates",
        15,
        "Arial",
    )
    shapes[9].width = 5_800_000

    replace_shape_text(
        shapes[14],
        "c) Estimate & summarize stratified seroincidence rates",
        15,
        "Arial",
    )
    replace_shape_text(
        shapes[13],
        "d) Visualize stratified seroincidence rates",
        15,
        "Arial",
    )
    replace_shape_text(
        shapes[19],
        "d) Visualize seroincidence estimates by stratum",
        15,
        "Arial",
    )

    plot = shapes[15]
    scale = 0.86
    plot.width = int(plot.width * scale)
    plot.height = int(plot.height * scale)


def main() -> None:
    prs = Presentation(SOURCE)
    revise_slide_one(prs)
    revise_slide_two(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)

    check = Presentation(OUTPUT)
    if len(check.slides) != 2:
        raise ValueError("Unexpected slide count after save")
    print(OUTPUT)


if __name__ == "__main__":
    main()
